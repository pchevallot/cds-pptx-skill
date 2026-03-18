"""
CdS PPTX Builder — Generateur de presentations PowerPoint
selon la charte graphique du Comptoir des Signaux.

Dependances :
    pip install python-pptx requests

Usage autonome :
    python cds_pptx.py

Usage en tant que module :
    from cds_pptx import CdsPptxBuilder

    builder = CdsPptxBuilder()
    builder.add_cover("Mon titre", "Sous-titre")
    builder.add_content_slide("Contexte", "Texte du contenu...")
    builder.add_table_slide("Resultats", ["Col1", "Col2"], [["a", "b"]])
    builder.save("output.pptx")
"""

import os
import tempfile
from io import BytesIO
from pathlib import Path

import requests
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

GITHUB_BASE = (
    "https://raw.githubusercontent.com/pchevallot/cds-pptx-skill/main/assets"
)

LOGO_URLS = {
    "bleu_jaune": f"{GITHUB_BASE}/logos/CDS-Logo-Bleu-Jaune.png",
    "jaune_blanc": f"{GITHUB_BASE}/logos/CDS-Logo-Jaune-Blanc.png",
    "bleu_blanc": f"{GITHUB_BASE}/logos/CDS-Logo-Bleu-Blanc.png",
    "noir": f"{GITHUB_BASE}/logos/CDS-Logo-Noir.png",
    "blanc": f"{GITHUB_BASE}/logos/CDS-Logo-Blanc.png",
}

MONOGRAMME_URLS = {
    "bleu_jaune": f"{GITHUB_BASE}/monogrammes/Monogramme-Bleu-Jaune.png",
    "blanc_jaune": f"{GITHUB_BASE}/monogrammes/Monogramme-Blanc-Jaune.png",
    "blanc": f"{GITHUB_BASE}/monogrammes/Monogramme-Blanc.png",
    "noir": f"{GITHUB_BASE}/monogrammes/Monogramme-Noir.png",
}

BANDEAU_URLS = {
    "bleu_h": f"{GITHUB_BASE}/bandeaux/Bandeau-motifs-bleu-horizontal.png",
    "blanc_h": f"{GITHUB_BASE}/bandeaux/Bandeau-motifs-blanc-horizontal.png",
    "jaune_h": f"{GITHUB_BASE}/bandeaux/Bandeau-motifs-jaune-horizontal.png",
    "bleu_v": f"{GITHUB_BASE}/bandeaux/Bandeau-motifs-bleu-vertical.png",
    "blanc_v": f"{GITHUB_BASE}/bandeaux/Bandeau-motifs-blanc-vertical.png",
    "jaune_v": f"{GITHUB_BASE}/bandeaux/Bandeau-motifs-jaune-vertical.png",
}

# ---------------------------------------------------------------------------
# Brand constants
# ---------------------------------------------------------------------------

CDS_BLEU = RGBColor(0x1F, 0x51, 0x9B)
CDS_OR = RGBColor(0xFD, 0xC9, 0x48)
CDS_BLANC = RGBColor(0xFF, 0xFF, 0xFF)
CDS_GRIS_FONCE = RGBColor(0x33, 0x33, 0x33)
CDS_GRIS_CLAIR = RGBColor(0xF5, 0xF5, 0xF5)

VERT = RGBColor(0x4C, 0xAF, 0x50)
ORANGE = RGBColor(0xFF, 0x98, 0x00)
ROUGE = RGBColor(0xF4, 0x43, 0x36)

FONT_NAME = "Open Sans"
FONT_FALLBACK = "Calibri"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ---------------------------------------------------------------------------
# Asset downloader with local cache
# ---------------------------------------------------------------------------

class AssetCache:
    """Downloads and caches remote assets to a local temp directory."""

    def __init__(self, cache_dir: str | Path | None = None):
        if cache_dir:
            self._dir = Path(cache_dir)
            self._dir.mkdir(parents=True, exist_ok=True)
        else:
            self._dir = Path(tempfile.mkdtemp(prefix="cds_pptx_"))

    def get(self, url: str) -> Path:
        """Download url to cache (if not already cached) and return local path."""
        filename = url.rsplit("/", 1)[-1]
        local = self._dir / filename
        if not local.exists():
            resp = requests.get(url, timeout=30)
            resp.raise_for_status()
            local.write_bytes(resp.content)
        return local


# ---------------------------------------------------------------------------
# Main builder
# ---------------------------------------------------------------------------

class CdsPptxBuilder:
    """
    Creates PPTX presentations following the CdS brand guide.

    All slides use:
    - 16:9 widescreen (13.333" x 7.5")
    - Open Sans font (fallback: Calibri)
    - CdS color palette (Bleu, Or, Blanc, Gris)
    - Logo placement per brand guidelines
    """

    def __init__(self, cache_dir: str | Path | None = None):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        self.cache = AssetCache(cache_dir)
        self._font = self._detect_font()

    def _detect_font(self) -> str:
        """Return the best available font name."""
        return FONT_NAME

    def _cover_title_size(self, title: str):
        """Return adaptive font size for cover/closing/section titles."""
        n = len(title)
        if n <= 30:
            return Pt(48)
        elif n <= 50:
            return Pt(40)
        elif n <= 80:
            return Pt(34)
        elif n <= 120:
            return Pt(28)
        return Pt(24)

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def add_cover(self, title: str, subtitle: str = "", date_str: str = ""):
        """
        Add a cover slide with blue background, centered logo, title and subtitle.

        The title font size adapts automatically to the text length so that
        long titles never overlap the subtitle or date.

        Args:
            title: Main title (e.g. "Bilan COPIL")
            subtitle: Subtitle (e.g. "Metropole du Lac Bleu")
            date_str: Date string (e.g. "18/03/2026"). Auto-generated if empty.
        """
        slide = self._blank_slide()

        # Blue background
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = CDS_BLEU

        # Logo (light version, centered — slightly smaller to leave room)
        self._add_centered_logo(slide, "jaune_blanc", height=Inches(1.0), top=Inches(0.5))

        # Title — adaptive font size, vertically centered in a generous area
        font_size = self._cover_title_size(title)
        self._add_textbox(
            slide,
            text=title,
            x=Inches(1), y=Inches(1.9), w=Inches(11.333), h=Inches(2.8),
            font_size=font_size, bold=True, color=CDS_BLANC,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
        )

        # Subtitle — positioned below the title area
        if subtitle:
            self._add_textbox(
                slide,
                text=subtitle,
                x=Inches(1), y=Inches(4.8), w=Inches(11.333), h=Inches(0.6),
                font_size=Pt(22), color=CDS_OR, align=PP_ALIGN.CENTER,
            )

        # Date
        if not date_str:
            from datetime import datetime
            date_str = datetime.now().strftime("%d/%m/%Y")
        self._add_textbox(
            slide,
            text=date_str,
            x=Inches(1), y=Inches(5.4), w=Inches(11.333), h=Inches(0.4),
            font_size=Pt(16), color=CDS_BLANC, align=PP_ALIGN.CENTER,
        )

        # Decorative pattern strip at bottom
        self._add_bandeau(slide, variant="jaune_h")

    def add_content_slide(self, title: str, content: str):
        """
        Add a content slide with blue title bar, logo, and text content.

        Args:
            title: Slide title (shown in the blue bar)
            content: Body text
        """
        slide = self._blank_slide()
        self._add_title_bar(slide, title)
        self._add_content_logo(slide)

        self._add_textbox(
            slide,
            text=content,
            x=Inches(0.5), y=Inches(1.5), w=Inches(12.333), h=Inches(5.5),
            font_size=Pt(16), color=CDS_GRIS_FONCE,
            align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, wrap=True,
        )

    def add_bullet_slide(self, title: str, bullets: list[str]):
        """
        Add a slide with a bulleted list.

        Args:
            title: Slide title
            bullets: List of bullet point strings
        """
        slide = self._blank_slide()
        self._add_title_bar(slide, title)
        self._add_content_logo(slide)

        textbox = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.2)
        )
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.font.name = self._font
            p.font.size = Pt(16)
            p.font.color.rgb = CDS_GRIS_FONCE
            p.level = 0
            p.space_after = Pt(12)

    def add_table_slide(
        self,
        title: str,
        headers: list[str],
        rows: list[list[str]],
        col_widths: list[float] | None = None,
    ):
        """
        Add a slide with a branded table.

        Args:
            title: Slide title
            headers: Column header texts
            rows: List of row data (each row is a list of strings)
            col_widths: Optional column widths in inches
        """
        slide = self._blank_slide()
        self._add_title_bar(slide, title)
        self._add_content_logo(slide)

        n_rows = len(rows) + 1  # +1 for header
        n_cols = len(headers)

        left = Inches(0.5)
        top = Inches(1.8)
        width = Inches(12.333)
        height = Inches(min(5.0, 0.4 * n_rows + 0.5))

        table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
        table = table_shape.table

        # Column widths
        if col_widths:
            for i, w in enumerate(col_widths):
                table.columns[i].width = Inches(w)
        else:
            col_w = 12.333 / n_cols
            for i in range(n_cols):
                table.columns[i].width = Inches(col_w)

        # Header row
        for col_idx, header_text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = CDS_BLEU
            p = cell.text_frame.paragraphs[0]
            p.font.name = self._font
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = CDS_BLANC
            p.alignment = PP_ALIGN.CENTER

        # Data rows
        for row_idx, row_data in enumerate(rows, start=1):
            for col_idx, cell_text in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_text)

                # Alternating row colors
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = CDS_GRIS_CLAIR

                p = cell.text_frame.paragraphs[0]
                p.font.name = self._font
                p.font.size = Pt(12)
                p.font.color.rgb = CDS_GRIS_FONCE

    def add_chart_slide(self, title: str, chart_image_path: str | Path):
        """
        Add a slide with a chart image (PNG/JPG).

        Args:
            title: Slide title
            chart_image_path: Path to the chart image file
        """
        slide = self._blank_slide()
        self._add_title_bar(slide, title)
        self._add_content_logo(slide)

        chart_path = Path(chart_image_path)
        if chart_path.exists():
            self._insert_chart_image(slide, chart_path)
        else:
            self._add_textbox(
                slide,
                text="Image non disponible",
                x=Inches(3), y=Inches(3.5), w=Inches(7.333), h=Inches(1),
                font_size=Pt(18), italic=True, color=CDS_GRIS_FONCE,
                align=PP_ALIGN.CENTER,
            )

    def add_radar_slide(
        self,
        title: str,
        labels: list[str],
        datasets: list[dict],
        chart_title: str = "",
    ):
        """
        Generate a radar chart with matplotlib and add it to the presentation.

        Uses CdS brand colors, square aspect ratio, and proper label placement.

        Args:
            title: Slide title (in the blue bar)
            labels: Axis labels (e.g. ["Axe 1", "Axe 2", ...])
            datasets: List of dicts with keys:
                - "label": Legend label
                - "values": List of numeric values (same length as labels)
                - "color": Optional hex color (default: CdS palette)
            chart_title: Optional chart title rendered inside the image
        """
        import math
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import numpy as np

        n = len(labels)
        angles = [i / n * 2 * math.pi for i in range(n)]
        angles += angles[:1]  # Close the polygon

        # CdS brand colors for datasets
        palette = ["#1F519B", "#FDC948", "#F44336", "#4CAF50", "#FF9800"]

        fig, ax = plt.subplots(figsize=(8, 8), subplot_kw=dict(polar=True))

        for i, ds in enumerate(datasets):
            values = ds["values"] + ds["values"][:1]
            color = ds.get("color", palette[i % len(palette)])
            ax.plot(angles, values, "o-", linewidth=2, label=ds["label"], color=color)
            ax.fill(angles, values, alpha=0.15, color=color)

        ax.set_xticks(angles[:-1])
        ax.set_xticklabels(labels, fontsize=11, fontfamily="sans-serif")

        # Ensure the radar is perfectly circular
        ax.set_aspect("equal")

        if chart_title:
            ax.set_title(
                chart_title, fontsize=16, fontweight="bold",
                color="#1F519B", pad=20, fontfamily="sans-serif",
            )

        ax.legend(
            loc="lower right", bbox_to_anchor=(1.15, -0.05),
            fontsize=10, framealpha=0.9,
        )

        # Save to temp file
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        fig.savefig(tmp.name, dpi=150, bbox_inches="tight", facecolor="white")
        plt.close(fig)

        # Add as chart slide
        slide = self._blank_slide()
        self._add_title_bar(slide, title)
        self._add_content_logo(slide)
        self._insert_chart_image(slide, Path(tmp.name))

        # Cleanup
        try:
            os.unlink(tmp.name)
        except OSError:
            pass

    def add_section_slide(self, title: str, subtitle: str = ""):
        """
        Add a section divider slide (blue background, centered text).

        Args:
            title: Section title
            subtitle: Optional subtitle
        """
        slide = self._blank_slide()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = CDS_BLEU

        font_size = self._cover_title_size(title)
        self._add_textbox(
            slide,
            text=title,
            x=Inches(1), y=Inches(2.0), w=Inches(11.333), h=Inches(2.5),
            font_size=font_size, bold=True, color=CDS_BLANC,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
        )

        if subtitle:
            self._add_textbox(
                slide,
                text=subtitle,
                x=Inches(1), y=Inches(4.7), w=Inches(11.333), h=Inches(1),
                font_size=Pt(20), color=CDS_OR, align=PP_ALIGN.CENTER,
            )

    def add_blocks_slide(self, title: str, blocks: list[dict]):
        """
        Add a slide with stacked content blocks, each with a colored vertical
        accent bar on the left — ideal for layered architectures, process steps,
        or categorized information.

        Args:
            title: Slide title (in the blue bar)
            blocks: List of dicts with keys:
                - "title": Block heading (displayed uppercase)
                - "content": Description text
                - "color": Optional hex color string (e.g. "#4CAF50").
                  Defaults cycle through CdS palette.
        """
        slide = self._blank_slide()
        self._add_title_bar(slide, title)
        self._add_content_logo(slide)

        palette = [CDS_BLEU, CDS_OR, VERT, ORANGE, ROUGE]
        n = len(blocks)
        if n == 0:
            return

        area_top = Inches(1.6)
        area_height = Inches(5.4)
        block_height = area_height // n
        gap = Inches(0.15)

        for i, block in enumerate(blocks):
            # Determine bar color
            color_val = block.get("color")
            if color_val and isinstance(color_val, str):
                h = color_val.lstrip("#")
                bar_color = RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
            else:
                bar_color = palette[i % len(palette)]

            block_top = area_top + block_height * i

            # Vertical accent bar
            bar_shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE.RECTANGLE
                Inches(0.7), block_top + gap,
                Inches(0.07), block_height - gap * 2,
            )
            bar_shape.fill.solid()
            bar_shape.fill.fore_color.rgb = bar_color
            bar_shape.line.fill.background()

            # Block title
            self._add_textbox(
                slide,
                text=block.get("title", "").upper(),
                x=Inches(1.0), y=block_top + gap,
                w=Inches(11), h=Inches(0.45),
                font_size=Pt(17), bold=True, color=bar_color,
            )

            # Block content
            content = block.get("content", "")
            if content:
                self._add_textbox(
                    slide,
                    text=content,
                    x=Inches(1.0), y=block_top + gap + Inches(0.5),
                    w=Inches(11), h=block_height - gap * 2 - Inches(0.55),
                    font_size=Pt(14), color=CDS_GRIS_FONCE,
                    valign=MSO_ANCHOR.TOP,
                )

    def add_cards_slide(
        self,
        title: str,
        cards: list[dict],
        footnote: str = "",
    ):
        """
        Add a slide with side-by-side cards (2-4), each with a colored top
        accent bar, a bold title, and description text — ideal for comparing
        concepts, pillars, or key facts.

        Args:
            title: Slide title (in the blue bar)
            cards: List of dicts (2-4 recommended) with keys:
                - "title": Card heading
                - "content": Card body text
                - "color": Optional hex color for the top bar
            footnote: Optional small text at the bottom of the slide
        """
        slide = self._blank_slide()
        self._add_title_bar(slide, title)
        self._add_content_logo(slide)

        palette = [CDS_OR, CDS_BLEU, VERT, ORANGE, ROUGE]
        n = len(cards)
        if n == 0:
            return

        margin_x = Inches(0.5)
        card_gap = Inches(0.3)
        total_width = SLIDE_WIDTH - margin_x * 2
        card_width = (total_width - card_gap * (n - 1)) // n
        card_top = Inches(1.6)
        card_height = Inches(4.6)
        bar_h = Inches(0.07)

        for i, card in enumerate(cards):
            card_left = margin_x + (card_width + card_gap) * i

            # Card background (light gray rounded rectangle)
            card_shape = slide.shapes.add_shape(
                5,  # MSO_SHAPE.ROUNDED_RECTANGLE
                card_left, card_top, card_width, card_height,
            )
            card_shape.fill.solid()
            card_shape.fill.fore_color.rgb = CDS_GRIS_CLAIR
            card_shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            card_shape.line.width = Pt(1)

            # Top accent bar
            color_val = card.get("color")
            if color_val and isinstance(color_val, str):
                h = color_val.lstrip("#")
                bar_color = RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
            else:
                bar_color = palette[i % len(palette)]

            bar = slide.shapes.add_shape(
                1,  # MSO_SHAPE.RECTANGLE
                card_left, card_top, card_width, bar_h,
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = bar_color
            bar.line.fill.background()

            # Card title
            self._add_textbox(
                slide,
                text=card.get("title", ""),
                x=card_left + Inches(0.25), y=card_top + Inches(0.4),
                w=card_width - Inches(0.5), h=Inches(0.5),
                font_size=Pt(16), bold=True, color=CDS_BLEU,
                align=PP_ALIGN.CENTER,
            )

            # Card content
            content = card.get("content", "")
            if content:
                self._add_textbox(
                    slide,
                    text=content,
                    x=card_left + Inches(0.25), y=card_top + Inches(1.0),
                    w=card_width - Inches(0.5), h=card_height - Inches(1.3),
                    font_size=Pt(12), color=CDS_GRIS_FONCE,
                    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.TOP,
                )

        # Footnote
        if footnote:
            self._add_textbox(
                slide,
                text=footnote,
                x=Inches(0.5), y=Inches(6.7), w=Inches(12.333), h=Inches(0.4),
                font_size=Pt(10), italic=True,
                color=RGBColor(0x99, 0x99, 0x99),
            )

    def add_closing_slide(self, text: str = "Merci de votre attention", contact: str = ""):
        """
        Add a closing slide with blue background.

        The text font size adapts automatically to avoid overlapping the contact
        information or the decorative bandeau.

        Args:
            text: Main closing message
            contact: Optional contact information
        """
        slide = self._blank_slide()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = CDS_BLEU

        self._add_centered_logo(slide, "jaune_blanc", height=Inches(1.0), top=Inches(0.5))

        # Closing text — adaptive font size
        font_size = self._cover_title_size(text)
        self._add_textbox(
            slide,
            text=text,
            x=Inches(1), y=Inches(1.9), w=Inches(11.333), h=Inches(2.2),
            font_size=font_size, bold=True, color=CDS_BLANC,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
        )

        if contact:
            self._add_textbox(
                slide,
                text=contact,
                x=Inches(1), y=Inches(4.2), w=Inches(11.333), h=Inches(1.5),
                font_size=Pt(16), color=CDS_OR, align=PP_ALIGN.CENTER,
            )

        # Decorative pattern strip at bottom
        self._add_bandeau(slide, variant="jaune_h")

    def save(self, path: str | Path) -> Path:
        """Save the presentation to a file."""
        path = Path(path)
        self.prs.save(str(path))
        return path

    def get_bytes(self) -> bytes:
        """Return the presentation as bytes (useful for web responses)."""
        buffer = BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------

    def _blank_slide(self):
        """Add a blank slide."""
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _add_title_bar(self, slide, title: str):
        """Add the standard blue title bar at the top of a content slide."""
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(0), Inches(0), SLIDE_WIDTH, Inches(1),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = CDS_BLEU
        shape.line.color.rgb = CDS_BLEU

        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self._font
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = CDS_BLANC
        p.alignment = PP_ALIGN.CENTER

    def _add_content_logo(self, slide):
        """Add the light logo to the top-right corner of a content slide (blue title bar)."""
        try:
            # Jaune-Blanc variant: yellow monogram + white text — visible on blue bar
            logo_path = self.cache.get(LOGO_URLS["jaune_blanc"])
            logo_h = Inches(0.5)
            logo_w = Inches(2.0)  # ~4:1 ratio
            left = SLIDE_WIDTH - logo_w - Inches(0.2)
            top = Inches(0.25)
            slide.shapes.add_picture(str(logo_path), left, top, height=logo_h)
        except Exception:
            pass  # Graceful degradation: no logo if download fails

    def _add_centered_logo(self, slide, variant: str, height, top):
        """Add a centered logo on cover/section slides."""
        try:
            url = LOGO_URLS.get(variant, LOGO_URLS["bleu_jaune"])
            logo_path = self.cache.get(url)
            logo_h = height
            logo_w = Inches(height / Inches(1) * 4)  # 4:1 ratio
            left = (SLIDE_WIDTH - logo_w) // 2
            slide.shapes.add_picture(str(logo_path), int(left), top, height=logo_h)
        except Exception:
            pass  # Graceful degradation

    def _insert_chart_image(self, slide, chart_path: Path):
        """Insert a chart image centered in the content area, preserving aspect ratio."""
        max_w = Inches(11.333)
        max_h = Inches(5.5)
        top = Inches(1.5)

        try:
            from PIL import Image as PILImage
            with PILImage.open(chart_path) as img:
                ratio = img.width / img.height
        except ImportError:
            ratio = 1.0  # Default to square (safe for radar charts)

        if ratio >= max_w / max_h:
            w = int(max_w)
            h = int(w / ratio)
        else:
            h = int(max_h)
            w = int(h * ratio)

        # Center horizontally, vertically center in content area
        left = (int(SLIDE_WIDTH) - w) // 2
        v_offset = (int(max_h) - h) // 2
        top = int(Inches(1.5)) + v_offset

        slide.shapes.add_picture(str(chart_path), left, top, width=w, height=h)

    def _add_bandeau(self, slide, variant="jaune_h", max_height_inches=1.2):
        """
        Add a decorative pattern strip at the bottom of a slide.

        The bandeau image is scaled to the full slide width while preserving
        its natural aspect ratio so the monogram patterns stay circular.

        Args:
            slide: The slide to add the bandeau to
            variant: Key from BANDEAU_URLS (default: yellow horizontal)
            max_height_inches: Maximum height cap in inches
        """
        try:
            url = BANDEAU_URLS.get(variant, BANDEAU_URLS["jaune_h"])
            bandeau_path = self.cache.get(url)

            # Get actual image dimensions to preserve aspect ratio
            try:
                from PIL import Image as PILImage
                with PILImage.open(bandeau_path) as img:
                    img_ratio = img.width / img.height
            except ImportError:
                img_ratio = 12.0  # Reasonable fallback for horizontal bandeau

            # Scale to full slide width, calculate height from ratio
            bandeau_w = int(SLIDE_WIDTH)
            bandeau_h = int(SLIDE_WIDTH / img_ratio)

            # Cap height so it doesn't take too much slide space
            max_h = Inches(max_height_inches)
            if bandeau_h > int(max_h):
                bandeau_h = int(max_h)
                bandeau_w = int(bandeau_h * img_ratio)

            left = (int(SLIDE_WIDTH) - bandeau_w) // 2
            top = int(SLIDE_HEIGHT) - bandeau_h

            slide.shapes.add_picture(
                str(bandeau_path), left, top,
                width=bandeau_w, height=bandeau_h,
            )
        except Exception:
            pass  # Graceful degradation

    def _add_textbox(
        self,
        slide,
        text: str,
        x, y, w, h,
        font_size=Pt(16),
        bold=False,
        italic=False,
        color=CDS_GRIS_FONCE,
        align=PP_ALIGN.LEFT,
        valign=None,
        wrap=True,
    ):
        """Add a formatted text box to a slide."""
        textbox = slide.shapes.add_textbox(x, y, w, h)
        tf = textbox.text_frame
        tf.word_wrap = wrap
        if valign:
            tf.vertical_anchor = valign

        p = tf.paragraphs[0]
        p.text = text
        p.font.name = self._font
        p.font.size = font_size
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
        p.alignment = align

        return textbox


# ---------------------------------------------------------------------------
# Standalone demo
# ---------------------------------------------------------------------------

def main():
    """Generate a demo presentation showcasing the CdS brand."""
    builder = CdsPptxBuilder()

    # Slide 1: Cover (long title to test adaptive sizing)
    builder.add_cover(
        title=(
            "Resilience des reseaux de telecommunication "
            "et datacenters a haute performance environnementale"
        ),
        subtitle="Enjeux, architectures et trajectoire pour les collectivites territoriales",
        date_str="Mars 2026",
    )

    # Slide 2: Section divider
    builder.add_section_slide(
        title="Contexte du projet",
        subtitle="Accompagnement a la transformation numerique",
    )

    # Slide 3: Cards (side-by-side comparison)
    builder.add_cards_slide(
        title="Contexte — Un ecosysteme sous tension",
        cards=[
            {
                "title": "Risques climatiques",
                "content": (
                    "Tempetes, canicules, inondations : "
                    "les infrastructures telecom subissent des stress croissants.\n\n"
                    "+40% d'incidents climatiques sur les reseaux depuis 2018 "
                    "(source : ARCEP, 2024)."
                ),
            },
            {
                "title": "Dependance numerique",
                "content": (
                    "93% des services publics territoriaux dependent "
                    "d'une connectivite reseau.\n\n"
                    "Une coupure de 4h = paralysie administrative "
                    "et rupture de service public."
                ),
            },
            {
                "title": "Empreinte carbone",
                "content": (
                    "Les datacenters representent 2,5% de la consommation "
                    "electrique nationale.\n\n"
                    "Objectif DNUM / ARCEP : -25% d'empreinte carbone "
                    "du numerique d'ici 2030."
                ),
            },
        ],
        footnote="Sources : ARCEP, Rapport sur l'etat d'internet en France 2024 — DNUM, Feuille de route numerique responsable 2025",
    )

    # Slide 4: Blocks (layered architecture)
    builder.add_blocks_slide(
        title="Architecture tri-couche resiliente",
        blocks=[
            {
                "title": "Couche Infrastructure",
                "content": "Fibre multi-operateur · LoRaWAN mutualise · Points de mutualisation · Energie secourue 72h+",
            },
            {
                "title": "Couche Services & Donnees",
                "content": "Datacenter HQE local/regional · Cloud souverain (SecNumCloud) · Replication geo-distribuee · Chiffrement bout en bout",
                "color": "#FDC948",
            },
            {
                "title": "Couche Usages & Gouvernance",
                "content": "PCA/PRA territorial · SOC mutualise · Supervision unifiee · Charte numerique responsable",
                "color": "#4CAF50",
            },
        ],
    )

    # Slide 5: Content (simple text)
    builder.add_content_slide(
        title="Contexte et enjeux",
        content=(
            "La Metropole du Lac Bleu a engage une demarche de transformation "
            "numerique et d'integration de l'intelligence artificielle dans ses "
            "services. Le Comptoir des Signaux accompagne cette transformation "
            "dans le cadre d'une mission d'AMO de 4 ans.\n\n"
            "Les enjeux principaux portent sur l'acculturation des agents, "
            "la securisation des donnees et la mise en place d'une gouvernance "
            "IA adaptee au contexte territorial."
        ),
    )

    # Slide 6: Bullet points
    builder.add_bullet_slide(
        title="Prochaines etapes",
        bullets=[
            "Finaliser la cartographie des cas d'usage prioritaires",
            "Deployer les pilotes sur 3 directions test",
            "Former les 120 agents referents IA",
            "Mettre en place le comite de gouvernance IA",
            "Preparer le bilan a mi-parcours pour le COPIL de septembre",
        ],
    )

    # Slide 7: Table
    builder.add_table_slide(
        title="Avancement par etape",
        headers=["Etape", "Intitule", "Statut", "Completion"],
        rows=[
            ["M1", "Acculturer et cadrer", "Termine", "100%"],
            ["M2", "Cartographier et prioriser", "En cours", "60%"],
            ["M3", "Securiser donnees et conformite", "A venir", "0%"],
            ["M4", "Architectures IA souveraines", "A venir", "0%"],
            ["M5", "Experimentation maitrisee", "A venir", "0%"],
        ],
    )

    # Slide 8: Closing (long text to test adaptive sizing)
    builder.add_closing_slide(
        text="Construisons ensemble une infrastructure numerique resiliente et responsable",
        contact=(
            "Pascal CHEVALLOT — Directeur de mission\n"
            "pchevallot@comptoirdessignaux.com — 06 02 03 40 13\n"
            "Le Comptoir des Signaux — comptoirdessignaux.com\n"
            "Architecte de trajectoire numerique publique integree"
        ),
    )

    output = builder.save("demo_cds_presentation.pptx")
    print(f"Presentation generee : {output}")
    print(f"Nombre de slides : {len(builder.prs.slides)}")


if __name__ == "__main__":
    main()
